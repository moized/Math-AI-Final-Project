from pptx import Presentation
from pptx.util import Inches, Pt

# Presentation setup
prs = Presentation()


def add_title_slide(title, subtitle):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullets_slide(title, bullets):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        p = body.add_paragraph() if i > 0 else body.paragraphs[0]
        p.text = b
        p.level = 0


def add_two_column_slide(title, left_items, right_items):
    slide_layout = prs.slide_layouts[3]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    left = slide.shapes.placeholders[1].text_frame
    left.clear()
    for i, b in enumerate(left_items):
        p = left.add_paragraph() if i > 0 else left.paragraphs[0]
        p.text = b
        p.level = 0
    right = slide.shapes.placeholders[2].text_frame
    right.clear()
    for i, b in enumerate(right_items):
        p = right.add_paragraph() if i > 0 else right.paragraphs[0]
        p.text = b
        p.level = 0


# Slides
add_title_slide(
    "Final Projesi: VAE ile Latent Manifold Analizi",
    "Ders: Mathematical Foundations of AI\nÖğrenci: Mohammed Izedin Mohammed\nÖğrenci No: STUDENT_ID_REMOVED"
)

add_bullets_slide(
    "Motivasyon ve Amaç",
    [
        "Yüksek boyutlu veriyi düşük boyutlu latent ölçü uzayına indirgeme",
        "VAE ile veri manifoldunun topolojisini öğrenme",
        "ELBO optimizasyonu ile olasılık modelleme"
    ]
)

add_bullets_slide(
    "Teorik Çerçeve (Measure Theory Bağlantısı)",
    [
        "Olasılık uzayı: (Ω, 𝔽, P) ve ölçülebilirlik",
        "ELBO: E_q[log p(x|z)] − D_KL(q(z|x) || p(z))",
        "Jensen Eşitsizliği → log p(x) için alt sınır",
        "Monte Carlo integrasyonu: örnekleme ile integral yaklaşımı"
    ]
)

add_bullets_slide(
    "Model Mimarisi (VAE)",
    [
        "Encoder: x → (μ, log σ²)",
        "Reparameterization: z = μ + σ ⊙ ε,  ε ~ N(0, I)",
        "Decoder: z → x̂",
        "Latent uzay: 2 boyutlu ölçü uzayı / manifold"
    ]
)

add_bullets_slide(
    "Kayıp Fonksiyonu (ELBO)",
    [
        "Reconstruction: Binary Cross-Entropy",
        "Regularization: KL Divergence",
        "Gaussian prior seçimi: KL için kapalı form çözüm",
        "Amaç: ELBO’yu maksimize ederek log-olabilirliği dolaylı artırmak"
    ]
)

add_bullets_slide(
    "Deneysel Kurulum",
    [
        "Veri: MNIST (60k eğitim / 10k test)",
        "Eğitim: Adam, 10 epoch",
        "Latent boyut: 2",
        "Ek analizler: latent grid, hata histogramı, β‑VAE ablation"
    ]
)

add_bullets_slide(
    "Sonuçlar: Latent Manifold",
    [
        "Sınıflar latent uzayda topolojik kümeler oluşturdu",
        "Manifold yapısı görsel olarak ortaya çıktı",
        "Latent grid haritası ile decoder davranışı gözlemlendi"
    ]
)

add_bullets_slide(
    "Sonuçlar: Optimizasyon ve Hata",
    [
        "ELBO kaybı epoch boyunca yakınsadı",
        "Reconstruction error histogramı: modelin zorlandığı örnekler",
        "β‑VAE karşılaştırması: KL ağırlığı temsil gücünü etkiliyor"
    ]
)

add_bullets_slide(
    "Örnek Üretimi ve Rekonstrüksiyon",
    [
        "Rastgele z örneklerinden sentetik rakam üretimi",
        "Orijinal vs rekonstrüksiyon karşılaştırması",
        "Latent temsillerin semantik bilgi taşıdığı gözlemlendi"
    ]
)

add_bullets_slide(
    "Sonuç ve Katkılar",
    [
        "Teori + uygulama birleşimi: ölçü uzayı, ELBO, KL",
        "Manifold öğrenimi görsel olarak doğrulandı",
        "Final katkısı: latent grid + hata analizi + β‑VAE"
    ]
)

add_bullets_slide(
    "Teşekkürler",
    [
        "Sorular?"
    ]
)

# Save
output_path = "VAE_Final_Project_Presentation_STUDENT_ID_REMOVED.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
